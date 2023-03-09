from pptx import Presentation

# 打开演示文稿
pr = Presentation(r"C:\Users\CN0CHHG\Documents\SRP\SRP\SRP_APAC_Overall 202302wo.pptx")

# 要删除的页码列表
pages_to_remove = list(range(2, 7)) + list(range(19, 24))+list(range(32, 37))+list(range(44, 49))+list(range(56, 61))+list(range(69, 79))

# 循环遍历每一页
for i, slide in enumerate(pr.slides):
    # 如果这一页是要删除的页码之一
    if i+1 in pages_to_remove:
        # 遍历每个形状
        for shape in slide.shapes:
            # 如果形状是图片类型
            if shape.shape_type == 13:
                # 从 XML 树中删除形状
                slide.shapes._spTree.remove(shape.element)

# 保存演示文稿
pr.save(r"C:\Users\CN0CHHG\Documents\SRP\SRP\SRP_APAC_Overall 202302wo.pptx")
print ("end")